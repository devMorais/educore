"""
EduCore PPTX Service — python-pptx native implementation (BS-011 / BS-014).

Design System:
  Primary Blue   #1A73E8
  Accent Orange  #FF8C00
  Dark           #212121
  Light          #FAFAFA

Slide layouts:
  cover | objectives | content_bullets | two_column |
  section_divider | quote_highlight | assessment |
  summary | closing  (+ aliases: stats_numbers, timeline)
"""

import os
import logging
import threading
from datetime import date
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# ── Palette ───────────────────────────────────────────────────────────────────
BLUE     = RGBColor(0x1A, 0x73, 0xE8)
BLUEDARK = RGBColor(0x0D, 0x47, 0xA1)
ORANGE   = RGBColor(0xFF, 0x8C, 0x00)
DARK     = RGBColor(0x21, 0x21, 0x21)
LIGHT    = RGBColor(0xFA, 0xFA, 0xFA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MIDGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MID = RGBColor(0x55, 0x55, 0x55)
GRAY_LIGHT_TEXT = RGBColor(0xBB, 0xDE, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _accent_rgb(accent: str | None) -> RGBColor:
    return {
        "orange": ORANGE,
        "blue":   BLUE,
        "green":  RGBColor(0x2E, 0x7D, 0x32),
        "purple": RGBColor(0x6A, 0x1B, 0x9A),
    }.get(accent or "blue", BLUE)


def _rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _oval(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _txb(slide, text: str, left, top, width, height, *,
         size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         font="Calibri", italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _bullet_list(slide, items: List[str], left, top, width, height,
                 size=17, color=DARK, bullet="▸  "):
    if not items:
        return
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.name  = "Calibri"
        run.font.size  = Pt(size)
        run.font.color.rgb = color


# ── Layout builders ───────────────────────────────────────────────────────────

def _cover(prs, slide_data: dict, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BLUEDARK)

    # Orange diagonal accent top-right
    _rect(slide, SLIDE_W - Inches(3.8), Emu(0), Inches(3.8), Inches(1.6), ORANGE)

    # EduCore brand label
    _txb(slide, "EduCore — IA Educacional",
         Inches(0.6), Inches(0.35), Inches(8), Inches(0.5),
         size=11, color=GRAY_LIGHT_TEXT)

    # Main title
    main_title = slide_data.get("title") or title
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(11.2), Inches(3.0))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = main_title
    run.font.name  = "Calibri"
    run.font.size  = Pt(38)
    run.font.bold  = True
    run.font.color.rgb = WHITE

    # Subtitle from first content item
    content = slide_data.get("content", [])
    if content:
        _txb(slide, content[0],
             Inches(0.6), Inches(4.7), Inches(9.5), Inches(0.7),
             size=18, color=GRAY_LIGHT_TEXT)

    # Date
    _txb(slide, date.today().strftime("%d/%m/%Y"),
         Inches(0.6), Inches(5.6), Inches(4), Inches(0.4),
         size=12, color=RGBColor(0x90, 0xCA, 0xF9))

    # Bottom orange stripe
    _rect(slide, Emu(0), SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), ORANGE)
    return slide


def _objectives(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    _rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.35), BLUE)
    _txb(slide, slide_data.get("title", "Objetivos de Aprendizagem"),
         Inches(0.5), Inches(0.25), Inches(11.5), Inches(0.9),
         size=26, bold=True, color=WHITE, font="Calibri")

    _rect(slide, Emu(0), Inches(1.35), Inches(0.1), SLIDE_H - Inches(1.35), ORANGE)

    content = slide_data.get("content", [])
    y = Inches(1.65)
    for i, item in enumerate(content[:6], 1):
        _oval(slide, Inches(0.25), y, Inches(0.42), Inches(0.42), ORANGE)
        _txb(slide, str(i),
             Inches(0.25), y, Inches(0.42), Inches(0.42),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, item, Inches(0.85), y, Inches(11.8), Inches(0.6),
             size=17, color=DARK)
        y += Inches(0.75)

    return slide


def _content_bullets(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    accent = _accent_rgb(slide_data.get("accent_color"))
    _rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.07), accent)
    _rect(slide, Emu(0), Inches(0.07), Inches(0.07), SLIDE_H - Inches(0.07), accent)
    _rect(slide, Inches(0.07), Emu(0), SLIDE_W - Inches(0.07), Inches(1.3), MIDGRAY)

    _txb(slide, slide_data.get("title", ""),
         Inches(0.3), Inches(0.12), Inches(12.1), Inches(1.1),
         size=24, bold=True, color=DARK)

    _bullet_list(slide, slide_data.get("content", [])[:7],
                 Inches(0.45), Inches(1.5), Inches(12.0), Inches(5.6))
    return slide


def _two_column(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.25), BLUE)
    _txb(slide, slide_data.get("title", ""),
         Inches(0.45), Inches(0.18), Inches(11.5), Inches(0.9),
         size=24, bold=True, color=WHITE)

    content = slide_data.get("content", [])
    mid = max(1, len(content) // 2)

    _rect(slide, Inches(0.3), Inches(1.45), Inches(5.9), Inches(5.7), MIDGRAY)
    _bullet_list(slide, content[:mid],
                 Inches(0.5), Inches(1.65), Inches(5.5), Inches(5.3), size=15)

    accent = _accent_rgb(slide_data.get("accent_color"))
    _rect(slide, Inches(6.35), Inches(1.6), Inches(0.05), Inches(5.5), accent)

    _rect(slide, Inches(6.55), Inches(1.45), Inches(6.4), Inches(5.7), MIDGRAY)
    _bullet_list(slide, content[mid:],
                 Inches(6.75), Inches(1.65), Inches(6.0), Inches(5.3), size=15)
    return slide


def _section_divider(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _accent_rgb(slide_data.get("accent_color"))
    _set_bg(slide, accent)

    # Decorative circle (lighter shade) top-right
    lighter = RGBColor(
        min(accent.red + 45, 255),
        min(accent.green + 45, 255),
        min(accent.blue + 45, 255),
    )
    _oval(slide, Inches(9.3), Inches(0.8), Inches(4.8), Inches(4.8), lighter)

    _txb(slide, slide_data.get("title", ""),
         Inches(0.7), Inches(2.2), Inches(8.6), Inches(2.5),
         size=40, bold=True, color=WHITE)

    content = slide_data.get("content", [])
    if content:
        _txb(slide, content[0],
             Inches(0.7), Inches(5.1), Inches(8.5), Inches(0.9),
             size=18, color=WHITE)
    return slide


def _quote_highlight(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    _rect(slide, Emu(0), Emu(0), Inches(0.28), SLIDE_H, ORANGE)
    _rect(slide, Emu(0), SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), ORANGE)

    _txb(slide, slide_data.get("title", ""),
         Inches(0.65), Inches(0.6), Inches(11.0), Inches(0.6),
         size=13, bold=True, color=BLUE)

    _txb(slide, "“",
         Inches(0.55), Inches(0.4), Inches(2), Inches(1.4),
         size=80, color=ORANGE, bold=True)

    content = slide_data.get("content", [])
    quote  = content[0] if content else slide_data.get("title", "")
    author = content[1] if len(content) > 1 else ""

    _txb(slide, quote,
         Inches(0.75), Inches(1.5), Inches(11.5), Inches(4.2),
         size=21, italic=True, color=DARK)

    if author:
        _txb(slide, f"— {author}",
             Inches(0.75), Inches(5.9), Inches(10), Inches(0.55),
             size=14, color=GRAY_MID)
    return slide


def _assessment(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.3), ORANGE)
    _txb(slide, slide_data.get("title", "Verificação de Aprendizagem"),
         Inches(0.5), Inches(0.2), Inches(11.5), Inches(0.9),
         size=26, bold=True, color=WHITE)

    content = slide_data.get("content", [])
    y = Inches(1.65)
    for i, q in enumerate(content[:5], 1):
        _rect(slide, Inches(0.3), y + Inches(0.04), Inches(0.46), Inches(0.46), ORANGE)
        _txb(slide, str(i),
             Inches(0.3), y + Inches(0.04), Inches(0.46), Inches(0.46),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, q, Inches(0.95), y, Inches(11.8), Inches(0.65), size=16, color=DARK)
        y += Inches(0.87)
    return slide


def _summary_slide(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    _rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.3), BLUE)
    _txb(slide, slide_data.get("title", "Síntese"),
         Inches(0.5), Inches(0.22), Inches(11.5), Inches(0.9),
         size=26, bold=True, color=WHITE)

    content = slide_data.get("content", [])
    mid = len(content) // 2 + len(content) % 2
    for ci, col in enumerate([content[:mid], content[mid:]]):
        x = Inches(0.4 + ci * 6.5)
        y = Inches(1.5)
        for item in col:
            _rect(slide, x, y + Inches(0.1), Inches(0.1), Inches(0.35), ORANGE)
            _txb(slide, item, x + Inches(0.2), y, Inches(5.9), Inches(0.6), size=14, color=DARK)
            y += Inches(0.72)
    return slide


def _closing(prs, slide_data: dict, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK)

    _rect(slide, SLIDE_W - Inches(4.2), Emu(0), Inches(4.2), Inches(1.2), ORANGE)

    _txb(slide, slide_data.get("title", "Obrigado!"),
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(2.0),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    content = slide_data.get("content", [])
    if content:
        _txb(slide, content[0],
             Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.85),
             size=18, color=GRAY_LIGHT_TEXT, align=PP_ALIGN.CENTER)

    _txb(slide, "Gerado pelo EduCore • IA Educacional",
         Inches(1), Inches(5.55), Inches(11.3), Inches(0.6),
         size=13, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.CENTER)

    _rect(slide, Emu(0), SLIDE_H - Inches(0.2), SLIDE_W, Inches(0.2), BLUE)
    return slide


_LAYOUT_MAP = {
    "cover":           (_cover,           True),
    "objectives":      (_objectives,      False),
    "content_bullets": (_content_bullets, False),
    "two_column":      (_two_column,      False),
    "section_divider": (_section_divider, False),
    "quote_highlight": (_quote_highlight, False),
    "assessment":      (_assessment,      False),
    "summary":         (_summary_slide,   False),
    "closing":         (_closing,         True),
    "stats_numbers":   (_content_bullets, False),
    "timeline":        (_content_bullets, False),
}

# Layouts that already have their own footer / don't need a generic one
_NO_FOOTER = {"cover", "closing", "section_divider"}


def _add_footer(slide, slide_num: int, total: int):
    """Adds EduCore brand + slide number to the slide bottom."""
    _txb(slide, "EduCore",
         Inches(0.3), SLIDE_H - Inches(0.38), Inches(2.5), Inches(0.32),
         size=9, color=RGBColor(0xAA, 0xAA, 0xAA))
    _txb(slide, f"{slide_num} / {total}",
         SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.38), Inches(1.2), Inches(0.32),
         size=9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


# ── Service ───────────────────────────────────────────────────────────────────

class PptxService:
    def generate(self, title: str, slides: List[dict], output_path: str) -> str:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        total = len(slides)
        for idx, slide_data in enumerate(slides, 1):
            key = (slide_data.get("layout") or "content_bullets").lower()
            builder, needs_title = _LAYOUT_MAP.get(key, (_content_bullets, False))
            if needs_title:
                slide = builder(prs, slide_data, title)
            else:
                slide = builder(prs, slide_data)
            if key not in _NO_FOOTER and slide is not None:
                _add_footer(slide, idx, total)

        prs.save(output_path)
        logger.info("PPTX gerado (python-pptx): %s", output_path)

        # Auto-cleanup after 5 minutes
        threading.Timer(300, self._cleanup, args=[output_path]).start()
        return output_path

    @staticmethod
    def _cleanup(path: str):
        try:
            if os.path.exists(path):
                os.remove(path)
                logger.info("PPTX temporário removido: %s", path)
        except OSError:
            pass


pptx_service = PptxService()
